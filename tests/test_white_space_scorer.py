import pytest
from pptx import Presentation
from pptx.util import Inches

from aesthetic_code.scorer.white_space_scorer import MarginWhiteSpaceScorer


@pytest.fixture
def mock_slide():
    # Create a presentation and a slide for testing
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Empty slide layout

    # Add shapes to mock the slide layout with various positions
    slide.shapes.add_shape(
        1, Inches(1), Inches(1), Inches(1), Inches(1)
    )  # Near top-left
    slide.shapes.add_shape(
        1, Inches(6), Inches(4), Inches(1), Inches(1)
    )  # Near bottom-right

    return slide


@pytest.fixture
def scorer(mock_slide):
    # Set slide width and height (e.g., using standard PowerPoint size of 10x7.5 inches)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    return MarginWhiteSpaceScorer(
        mock_slide, slide_width, slide_height, measurement_unit="inches"
    )


def test_bounding_box(scorer):
    # Test if bounding box calculations are accurate
    bbox = scorer._get_bounding_box()
    assert bbox["left"] > 0, "Bounding box left position should be greater than zero"
    assert (
        bbox["right"] < scorer._width
    ), "Bounding box right position should be less than slide width"
    assert bbox["top"] > 0, "Bounding box top position should be greater than zero"
    assert (
        bbox["bottom"] < scorer._height
    ), "Bounding box bottom position should be less than slide height"


def test_margin_calculation(scorer):
    # Test if margin calculations are within expected range
    margins = scorer._calculate_margins()
    assert (
        "left" in margins and margins["left"] >= 0
    ), "Left margin should be calculated and non-negative"
    assert (
        "right" in margins and margins["right"] >= 0
    ), "Right margin should be calculated and non-negative"
    assert (
        "top" in margins and margins["top"] >= 0
    ), "Top margin should be calculated and non-negative"
    assert (
        "bottom" in margins and margins["bottom"] >= 0
    ), "Bottom margin should be calculated and non-negative"


def test_white_space_score(scorer):
    # Test if the white space score is calculated correctly
    score = scorer.calculate_white_space_score()
    assert 0.0 <= score <= 1.0, "Whitespace score should be a float between 0 and 1"

    # Based on the positions of shapes, expect a certain score
    # (e.g., if they satisfy horizontal/vertical margin thresholds partially)
    # Adjust expected_score based on actual margin settings and thresholds
    expected_score = 0.5  # Example value based on expected margin setup
    assert score == pytest.approx(
        expected_score, 0.1
    ), f"Expected whitespace score around {expected_score}, got {score}"
