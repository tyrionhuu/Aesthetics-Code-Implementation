import pytest
from pptx import Presentation
from pptx.util import Inches

from aesthetic_code.scorer.font_hierarchy_scorer import FontHierarchyScorer


@pytest.fixture
def mock_slide():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    # set font size
    title.text_frame.paragraphs[0].font.size = Inches(1)
    subtitle.text_frame.paragraphs[0].font.size = Inches(0.5)

    return slide


@pytest.fixture
def scorer():
    scorer = FontHierarchyScorer()
    return scorer


def test_font_hierarchy_scorer(scorer, mock_slide):
    scorer._title_node = mock_slide.shapes.title
    scorer._subtitle_node = mock_slide.placeholders[1]
    assert scorer.score() == 1.0
