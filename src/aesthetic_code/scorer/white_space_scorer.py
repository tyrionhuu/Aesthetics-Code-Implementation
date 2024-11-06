from typing import TypeAlias, Union

from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import BasePlaceholder
from pptx.util import Length

from aesthetic_code.segmenter.segmenter import SegmentTreeNode
from aesthetic_code.utils import unit_conversion

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
    BasePlaceholder,
]


class MarginWhiteSpaceScorer:
    """
    This class is responsible for scoring the white space in the margins of a slide.
    """

    def __init__(
        self,
        # shapes: list[Shape],
        segment_tree: SegmentTreeNode,
        slide_width: Length,
        slide_height: Length,
        measurement_unit: str = "pt",
    ):
        self._segment_tree = segment_tree
        self._shapes = segment_tree.shapes
        self._measurement_unit = measurement_unit
        self._width = unit_conversion(slide_width, self._measurement_unit)
        self._height = unit_conversion(slide_height, self._measurement_unit)
        self._margin_threshold = {
            "horizontal": (0.1 * self._width, 0.3 * self._width),
            "vertical": (0.1 * self._height, 0.3 * self._height),
        }

    def _get_bounding_box(self) -> dict:
        leftmost_x = self._width
        rightmost_x = 0.0
        topmost_y = self._height
        bottommost_y = 0.0
        for shape in self._shapes:
            left = unit_conversion(shape.left, self._measurement_unit)
            top = unit_conversion(shape.top, self._measurement_unit)
            width = unit_conversion(shape.width, self._measurement_unit)
            height = unit_conversion(shape.height, self._measurement_unit)
            right = left + width
            bottom = top + height
            leftmost_x = min(leftmost_x, left)
            rightmost_x = max(rightmost_x, right)
            topmost_y = min(topmost_y, top)
            bottommost_y = max(bottommost_y, bottom)

        return {
            "left": leftmost_x,
            "right": rightmost_x,
            "top": topmost_y,
            "bottom": bottommost_y,
        }

    def _calculate_margins(self) -> dict:
        bounding_box = self._get_bounding_box()
        left_margin = bounding_box["left"]
        right_margin = self._width - bounding_box["right"]
        top_margin = bounding_box["top"]
        bottom_margin = self._height - bounding_box["bottom"]
        return {
            "left": left_margin,
            "right": right_margin,
            "top": top_margin,
            "bottom": bottom_margin,
        }

    def calculate_white_space_score(self) -> float:
        margins = self._calculate_margins()
        horizontal_margin_score = 0
        vertical_margin_score = 0
        if (
            margins["left"] >= self._margin_threshold["horizontal"][0]
            and margins["left"] <= self._margin_threshold["horizontal"][1]
        ):
            horizontal_margin_score += 1
        if (
            margins["right"] >= self._margin_threshold["horizontal"][0]
            and margins["right"] <= self._margin_threshold["horizontal"][1]
        ):
            horizontal_margin_score += 1
        if (
            margins["top"] >= self._margin_threshold["vertical"][0]
            and margins["top"] <= self._margin_threshold["vertical"][1]
        ):
            vertical_margin_score += 1
        if (
            margins["bottom"] >= self._margin_threshold["vertical"][0]
            and margins["bottom"] <= self._margin_threshold["vertical"][1]
        ):
            vertical_margin_score += 1
        return float((horizontal_margin_score + vertical_margin_score) / 4)


class GroupSpacingWhiteSpaceScorer:
    """
    This class is responsible for scoring the white space between groups of shapes on a slide.
    It evaluates the vertical and horizontal spacing between groups of elements and assigns a score.
    """

    def __init__(
        self,
        segment_tree: SegmentTreeNode,
    ):
        self._segment_tree = segment_tree

    def score(self) -> float:
        """
        Calculate the overall white space score for the segment tree.
        The score is based on the spacing between the groups of shapes.
        """
        return self._score_subregions(self._segment_tree)

    def _score_subregions(self, node: SegmentTreeNode) -> float:
        """
        Recursively score the white space for each subregion of the segment tree.
        The score is calculated based on the distances between groups of shapes.
        """
        if node.is_leaf():
            return 0.0  # No space between groups in a leaf node

        total_score = 0.0
        subregion_count = len(node.subregions)

        for i, subregion in enumerate(node.subregions):
            # Score the distance between the current subregion and the next one
            if i < subregion_count - 1:
                next_subregion = node.subregions[i + 1]
                total_score += self._score_spacing_between_groups(
                    subregion, next_subregion
                )

            # Recursively score subregions within the current subregion
            total_score += self._score_subregions(subregion)

        return total_score

    def _score_spacing_between_groups(
        self, group1: SegmentTreeNode, group2: SegmentTreeNode
    ) -> float:
        """
        Calculate the score for the spacing between two groups of shapes.
        """
        # Calculate the minimum vertical and horizontal distance between groups
        vertical_distance = self._calculate_vertical_distance(group1, group2)
        horizontal_distance = self._calculate_horizontal_distance(group1, group2)

        # Combine the distances into a single score, with weights for vertical and horizontal spacing
        return vertical_distance * 0.7 + horizontal_distance * 0.3

    def _calculate_vertical_distance(
        self, group1: SegmentTreeNode, group2: SegmentTreeNode
    ) -> float:
        """
        Calculate the minimum vertical distance between two groups of shapes.
        """
        top1 = min([shape.top for shape in group1.shapes])
        bottom1 = max([shape.top + shape.height for shape in group1.shapes])
        top2 = min([shape.top for shape in group2.shapes])
        bottom2 = max([shape.top + shape.height for shape in group2.shapes])

        # Calculate vertical distance as the gap between the bottom of one group and the top of the other
        return max(0.0, min(top1, top2) - max(bottom1, bottom2))

    def _calculate_horizontal_distance(
        self, group1: SegmentTreeNode, group2: SegmentTreeNode
    ) -> float:
        """
        Calculate the minimum horizontal distance between two groups of shapes.
        """
        left1 = min([shape.left for shape in group1.shapes])
        right1 = max([shape.left + shape.width for shape in group1.shapes])
        left2 = min([shape.left for shape in group2.shapes])
        right2 = max([shape.left + shape.width for shape in group2.shapes])

        # Calculate horizontal distance as the gap between the right of one group and the left of the other
        return max(0.0, min(left1, left2) - max(right1, right2))
