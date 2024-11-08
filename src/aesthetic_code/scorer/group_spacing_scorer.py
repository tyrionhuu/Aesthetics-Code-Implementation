from typing import TypeAlias, Union, cast

from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import BasePlaceholder
from pptx.util import Length

from aesthetic_code.segmenter.segmenter import SegmentTreeNode, get_all_neighbor_pairs
from aesthetic_code.utils import unit_conversion

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
    BasePlaceholder,
]

Subregion: TypeAlias = Union[Shape, "SegmentTreeNode"]


class GroupSpacingScorer:
    """
    This class is responsible for scoring the white space between groups of shapes on a slide.
    It evaluates the vertical and horizontal spacing between groups of elements and assigns a score.
    """

    def __init__(
        self,
        slide_width: Length,
        slide_height: Length,
        segment_tree: SegmentTreeNode,
        spacing_threshold: list[float] = [0.1, 0.3],
        unit_measurement: str = "pt",
    ):
        self._segment_tree = segment_tree
        self._neighbor_pairs = get_all_neighbor_pairs(segment_tree)
        self._spacing_threshold = spacing_threshold
        self._unit_measurement = unit_measurement
        self._slide_width = slide_width
        self._slide_height = slide_height

    @property
    def spacing_threshold(self) -> list[float]:
        return self._spacing_threshold

    @spacing_threshold.setter
    def spacing_threshold(self, value: list[float]):
        self._spacing_threshold = value

    @property
    def unit_measurement(self) -> str:
        return self._unit_measurement

    @unit_measurement.setter
    def unit_measurement(self, value: str):
        self._unit_measurement = value

    @property
    def slide_width(self) -> Length:
        return self._slide_width

    @slide_width.setter
    def slide_width(self, value: Length):
        self._slide_width = value

    @property
    def slide_height(self) -> Length:
        return self._slide_height

    @slide_height.setter
    def slide_height(self, value: Length):
        self._slide_height = value

    def score(self) -> float:
        """
        Calculate the overall white space score for the segment tree.
        The score is based on the spacing between the groups of shapes.
        """
        scores = []
        for pair in self._neighbor_pairs:
            scores.append(self._score_pair(pair))

        return sum(scores) / len(scores)

    def _score_pair(self, pair: tuple[str, Subregion, Subregion]) -> float:
        """
        Calculate the white space score between two subregions.
        The score is based on the spacing between the two subregions.
        """
        direction, subregion1, subregion2 = pair
        if direction == "vertical":
            return self._score_horizontal_spacing(subregion1, subregion2)
        elif direction == "horizontal":
            return self._score_vertical_spacing(subregion1, subregion2)
        else:
            return self._score_segment_spacing(subregion1, subregion2)

    def _score_segment_spacing(
        self, subregion1: Subregion, subregion2: Subregion
    ) -> float:
        """
        Calculate the white space score between two subregions that are not directly adjacent.
        The score is based on the spacing between the two subregions.
        """

        assert isinstance(subregion2, SegmentTreeNode)
        if isinstance(subregion1, SegmentTreeNode):
            subregion = cast(SegmentTreeNode, subregion1)
            right_subregion1 = subregion.bounding_box["right"]
            left_subregion1 = subregion.bounding_box["left"]
            top_subregion1 = subregion.bounding_box["top"]
            bottom_subregion1 = subregion.bounding_box["bottom"]
        else:
            shape = cast(Shape, subregion1)
            right_subregion1 = unit_conversion(
                shape.left, self._unit_measurement
            ) + unit_conversion(shape.width, self._unit_measurement)
            left_subregion1 = unit_conversion(shape.left, self._unit_measurement)
            top_subregion1 = unit_conversion(shape.top, self._unit_measurement)
            bottom_subregion1 = unit_conversion(
                shape.top, self._unit_measurement
            ) + unit_conversion(shape.height, self._unit_measurement)

        horizontal_spacing = 0.0

        if right_subregion1 <= subregion2.bounding_box["left"]:
            horizontal_spacing = subregion2.bounding_box["left"] - right_subregion1
        elif subregion2.bounding_box["right"] <= left_subregion1:
            horizontal_spacing = left_subregion1 - subregion2.bounding_box["right"]
        else:
            pass

        vertical_spacing = 0.0

        if bottom_subregion1 <= subregion2.bounding_box["top"]:
            vertical_spacing = subregion2.bounding_box["top"] - bottom_subregion1
        elif subregion2.bounding_box["bottom"] <= top_subregion1:
            vertical_spacing = top_subregion1 - subregion2.bounding_box["bottom"]
        else:
            pass

        score = 0.0

        if horizontal_spacing <= self._spacing_threshold[0] * unit_conversion(
            self._slide_width, self._unit_measurement
        ) or horizontal_spacing >= self._spacing_threshold[1] * unit_conversion(
            self._slide_width, self._unit_measurement
        ):
            score += 0.0
        else:
            score += 0.5

        if vertical_spacing <= self._spacing_threshold[0] * unit_conversion(
            self._slide_height, self._unit_measurement
        ) or vertical_spacing >= self._spacing_threshold[1] * unit_conversion(
            self._slide_height, self._unit_measurement
        ):
            score += 0.0
        else:
            score += 0.5

        return score

    def _score_horizontal_spacing(
        self, subregion1: Subregion, subregion2: Subregion
    ) -> float:
        """
        Calculate the white space score between two subregions in the horizontal direction.
        The score is based on the horizontal spacing between the two subregions.
        """
        if isinstance(subregion1, SegmentTreeNode):
            subregion = cast(SegmentTreeNode, subregion1)
            right_subregion1 = subregion.bounding_box["right"]
            left_subregion1 = subregion.bounding_box["left"]
        else:
            shape = cast(Shape, subregion1)
            right_subregion1 = unit_conversion(
                shape.left, self._unit_measurement
            ) + unit_conversion(shape.width, self._unit_measurement)
            left_subregion1 = unit_conversion(shape.left, self._unit_measurement)

        if isinstance(subregion2, SegmentTreeNode):
            subregion = cast(SegmentTreeNode, subregion2)
            right_subregion2 = subregion.bounding_box["right"]
            left_subregion2 = subregion.bounding_box["left"]
        else:
            shape = cast(Shape, subregion2)
            right_subregion2 = unit_conversion(
                shape.left, self._unit_measurement
            ) + unit_conversion(shape.width, self._unit_measurement)
            left_subregion2 = unit_conversion(shape.left, self._unit_measurement)

        if right_subregion1 <= left_subregion2:
            spacing = left_subregion2 - right_subregion1
        elif right_subregion2 <= left_subregion1:
            spacing = left_subregion1 - right_subregion2
        else:
            raise ValueError("Overlapping subregions in horizontal spacing calculation")

        if spacing <= self._spacing_threshold[0] * unit_conversion(
            self._slide_width, self._unit_measurement
        ) or spacing >= self._spacing_threshold[1] * unit_conversion(
            self._slide_width, self._unit_measurement
        ):
            return 0.0
        else:
            return 1.0

    def _score_vertical_spacing(
        self, subregion1: Subregion, subregion2: Subregion
    ) -> float:
        """
        Calculate the white space score between two subregions in the vertical direction.
        The score is based on the vertical spacing between the two subregions.
        """
        if isinstance(subregion1, SegmentTreeNode):
            subregion = cast(SegmentTreeNode, subregion1)
            bottom_subregion1 = subregion.bounding_box["bottom"]
            top_subregion1 = subregion.bounding_box["top"]
        else:
            shape = cast(Shape, subregion1)
            bottom_subregion1 = unit_conversion(
                shape.top, self._unit_measurement
            ) + unit_conversion(shape.height, self._unit_measurement)
            top_subregion1 = unit_conversion(shape.top, self._unit_measurement)

        if isinstance(subregion2, SegmentTreeNode):
            subregion = cast(SegmentTreeNode, subregion2)
            bottom_subregion2 = subregion.bounding_box["bottom"]
            top_subregion2 = subregion.bounding_box["top"]
        else:
            shape = cast(Shape, subregion2)
            bottom_subregion2 = unit_conversion(
                shape.top, self._unit_measurement
            ) + unit_conversion(shape.height, self._unit_measurement)
            top_subregion2 = unit_conversion(shape.top, self._unit_measurement)

        if bottom_subregion1 <= top_subregion2:
            spacing = top_subregion2 - bottom_subregion1
        elif bottom_subregion2 <= top_subregion1:
            spacing = top_subregion1 - bottom_subregion2
        else:
            raise ValueError("Overlapping subregions in vertical spacing calculation")

        if spacing <= self._spacing_threshold[0] * unit_conversion(
            self._slide_height, self._unit_measurement
        ) or spacing >= self._spacing_threshold[1] * unit_conversion(
            self._slide_height, self._unit_measurement
        ):
            return 0.0
        else:
            return 1.0
