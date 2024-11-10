from typing import TypeAlias, Union

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import BasePlaceholder
from pptx.text.text import Font, TextFrame
from pptx.util import Length

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
    BasePlaceholder,
]


class FontHierarchyScorer:
    def __init__(self):
        self._title_node = None
        self._subtitle_node = None
        self._body_node = None

    @property
    def title_node(self):
        return self._title_node

    @title_node.setter
    def title_node(self, title_node):
        if not isinstance(title_node, Shape):  # type: ignore[misc]
            raise ValueError("title_node must be a Shape object")
        if title_node.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
            raise ValueError("title_node must be a placeholder shape")
        if title_node.placeholder_format != PP_PLACEHOLDER_TYPE.TITLE:
            raise ValueError("title_node must be a title placeholder")
        if not title_node.has_text_frame:
            raise ValueError("title_node must have a text frame")
        self._title_node = title_node

    @property
    def subtitle_node(self):
        return self._subtitle_node

    @subtitle_node.setter
    def subtitle_node(self, subtitle_node):
        if not isinstance(subtitle_node, Shape):  # type: ignore[misc]
            raise ValueError("subtitle_node must be a Shape object")
        if subtitle_node.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
            raise ValueError("subtitle_node must be a placeholder shape")
        if subtitle_node.placeholder_format != PP_PLACEHOLDER_TYPE.SUBTITLE:
            raise ValueError("subtitle_node must be a subtitle placeholder")
        if not subtitle_node.has_text_frame:
            raise ValueError("subtitle_node must have a text frame")
        self._subtitle_node = subtitle_node

    @property
    def body_node(self):
        return self._body_node

    @body_node.setter
    def body_node(self, body_node):
        if not isinstance(body_node, Shape):  # type: ignore[misc]
            raise ValueError("body_node must be a Shape object")
        if not body_node.has_text_frame:
            raise ValueError("body_node must have a text frame")
        self._body_node = body_node

    def _get_text_frame(self, shape: Shape) -> TextFrame:
        if not shape.has_text_frame:
            raise ValueError("Shape does not have a text frame")
        if hasattr(shape, "text_frame"):
            return shape.text_frame
        else:
            raise ValueError("Shape does not have text frame attribute")

    def _get_font_objects(self, text_frame: TextFrame) -> list[Font]:
        return [paragraph.font for paragraph in text_frame.paragraphs]

    def _get_font_sizes(self, font_objects: list[Font]) -> list[Length]:
        return [font.size for font in font_objects if font.size is not None]

    def score(self) -> float:
        return self._score_font_hierarchy()

    def _score_font_hierarchy(self) -> float:
        if all([self._title_node, self._subtitle_node, self._body_node]):
            title_text_frame = self._get_text_frame(self._title_node)
            subtitle_text_frame = self._get_text_frame(self._subtitle_node)
            body_text_frame = self._get_text_frame(self._body_node)

            title_font_objects = self._get_font_objects(title_text_frame)
            subtitle_font_objects = self._get_font_objects(subtitle_text_frame)
            body_font_objects = self._get_font_objects(body_text_frame)

            title_font_sizes = self._get_font_sizes(title_font_objects)
            subtitle_font_sizes = self._get_font_sizes(subtitle_font_objects)
            body_font_sizes = self._get_font_sizes(body_font_objects)

            title_font_size = sum(title_font_sizes) / len(title_font_sizes)
            subtitle_font_size = sum(subtitle_font_sizes) / len(subtitle_font_sizes)
            body_font_size = sum(body_font_sizes) / len(body_font_sizes)

            if title_font_size > subtitle_font_size > body_font_size:
                return 1.0
            return 0.0
        elif all([self._title_node, self._subtitle_node]):
            title_text_frame = self._get_text_frame(self._title_node)
            subtitle_text_frame = self._get_text_frame(self._subtitle_node)

            title_font_objects = self._get_font_objects(title_text_frame)
            subtitle_font_objects = self._get_font_objects(subtitle_text_frame)

            title_font_sizes = self._get_font_sizes(title_font_objects)
            subtitle_font_sizes = self._get_font_sizes(subtitle_font_objects)

            title_font_size = sum(title_font_sizes) / len(title_font_sizes)
            subtitle_font_size = sum(subtitle_font_sizes) / len(subtitle_font_sizes)

            if title_font_size > subtitle_font_size:
                return 1.0
            return 0.0
        elif all([self._title_node, self._body_node]):
            title_text_frame = self._get_text_frame(self._title_node)
            body_text_frame = self._get_text_frame(self._body_node)

            title_font_objects = self._get_font_objects(title_text_frame)
            body_font_objects = self._get_font_objects(body_text_frame)

            title_font_sizes = self._get_font_sizes(title_font_objects)
            body_font_sizes = self._get_font_sizes(body_font_objects)

            title_font_size = sum(title_font_sizes) / len(title_font_sizes)
            body_font_size = sum(body_font_sizes) / len(body_font_sizes)

            if title_font_size > body_font_size:
                return 1.0
            return 0.0
        elif all([self._subtitle_node, self._body_node]):
            subtitle_text_frame = self._get_text_frame(self._subtitle_node)
            body_text_frame = self._get_text_frame(self._body_node)

            subtitle_font_objects = self._get_font_objects(subtitle_text_frame)
            body_font_objects = self._get_font_objects(body_text_frame)

            subtitle_font_sizes = self._get_font_sizes(subtitle_font_objects)
            body_font_sizes = self._get_font_sizes(body_font_objects)

            subtitle_font_size = sum(subtitle_font_sizes) / len(subtitle_font_sizes)
            body_font_size = sum(body_font_sizes) / len(body_font_sizes)

            if subtitle_font_size > body_font_size:
                return 1.0
            return 0.0
        else:
            raise ValueError("Two of the three nodes must be provided")
