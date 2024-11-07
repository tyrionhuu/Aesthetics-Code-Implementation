from typing import TypeAlias, Union, cast

from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture
from pptx.shapes.placeholder import BasePlaceholder
from pptx.util import Length

from aesthetic_code.utils import intervals_minus_interval, unit_conversion

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
    BasePlaceholder,
]

Subregion: TypeAlias = Union[Shape, "SegmentTreeNode"]


class SegmentTreeNode:
    def __init__(
        self,
        direction: str = "leaf",  # Direction of arrangement, either 'vertical' or 'horizontal'
        subregions: list = [],  # Subregions of this node
        bounding_box: dict[str, float] = {},  # Bounding box of the node
    ):
        """
        Initialize a SegmentTreeNode with information about its subtree arrangement.

        Args:
            direction (str): Indicates how the node's children are arranged,
                             either 'vertical' or 'horizontal'.
            subregions (list[SegmentTreeNode]): list of subregions in this region.
        """
        self._subregions = subregions
        self._direction = direction  # Arrangement of subregions
        self._bounding_box = bounding_box  # Bounding box of the node

    @property
    def subregions(self) -> list[Subregion]:
        return self._subregions

    @property
    def bounding_box(self) -> dict[str, float]:
        return self._bounding_box

    @property
    def direction(self) -> str:
        return self._direction

    def is_leaf(self) -> bool:
        """
        Determines if the node is a leaf node (i.e., has no subregions).

        Returns:
            bool: True if the node is a leaf node, False otherwise.
        """
        return self._direction == "leaf"

    def print_tree(self, level: int = 0, indent: str = "  "):
        """
        Print a visual representation of the segment tree.

        Args:
            level (int): The current level of the tree (default is 0).
        """
        if self.is_leaf():
            if isinstance(self._subregions, list):
                shapes = cast(list[Shape], self._subregions)
            else:
                raise ValueError("Invalid subregion type")
            print(
                f"{indent * level}Leaf: "
                + ", ".join([str(shape.shape_type) for shape in shapes])
                + f"\t{self._bounding_box}"
            )
        else:
            print(f"{indent * level}{self._direction}:" + f"\t{self._bounding_box}")
            for subregion in self._subregions:
                if isinstance(subregion, SegmentTreeNode):
                    subregion.print_tree(level + 1)
                else:
                    raise ValueError("Invalid subregion type")


class Segmenter:
    def __init__(
        self,
        shapes: list[Shape],
        slide_width: Length | None,
        slide_height: Length | None,
        measurement_unit: str = "pt",
    ):
        self._shapes = shapes
        self._measurement_unit = measurement_unit
        self._slide_width = unit_conversion(slide_width, self._measurement_unit)
        self._slide_height = unit_conversion(slide_height, self._measurement_unit)

    def __call__(self, *args, **kwds):
        return self.segment()

    def segment(self) -> SegmentTreeNode:
        return self._segment_region(self._shapes)

    def _segment_region(self, shapes: list) -> SegmentTreeNode:
        if not shapes:  # If shapes list is empty, return a leaf node
            raise ValueError("No shapes to segment")

        if len(shapes) == 1:
            shape = shapes[0]
            bounding_box = self._generate_bounding_box_from_shapes(shape)
            return SegmentTreeNode(
                direction="leaf", subregions=shapes, bounding_box=bounding_box
            )

        split_directions = ["horizontal", "vertical"]
        for split_direction in split_directions:
            subregions = self._try_split(shapes, split_direction)
            if subregions:
                child_nodes = [
                    self._segment_region(subregion) for subregion in subregions
                ]
                bounding_boxs = [node._bounding_box for node in child_nodes]
                leftmost = min([box["left"] for box in bounding_boxs])
                topmost = min([box["top"] for box in bounding_boxs])
                rightmost = max([box["right"] for box in bounding_boxs])
                bottommost = max([box["bottom"] for box in bounding_boxs])
                return SegmentTreeNode(
                    direction=split_direction,
                    subregions=child_nodes,
                    bounding_box={
                        "left": leftmost,
                        "top": topmost,
                        "right": rightmost,
                        "bottom": bottommost,
                    },
                )
        bounding_boxs = [
            self._generate_bounding_box_from_shapes(shape) for shape in shapes
        ]
        leftmost = min([box["left"] for box in bounding_boxs])
        topmost = min([box["top"] for box in bounding_boxs])
        rightmost = max([box["right"] for box in bounding_boxs])
        bottommost = max([box["bottom"] for box in bounding_boxs])
        return SegmentTreeNode(
            direction="leaf",
            subregions=shapes,
            bounding_box={
                "left": leftmost,
                "top": topmost,
                "right": rightmost,
                "bottom": bottommost,
            },
        )

    def _generate_bounding_box_from_shapes(self, shape: Shape) -> dict:
        return {
            "left": unit_conversion(shape.left, self._measurement_unit),
            "top": unit_conversion(shape.top, self._measurement_unit),
            "right": unit_conversion(shape.left, self._measurement_unit)
            + unit_conversion(shape.width, self._measurement_unit),
            "bottom": unit_conversion(shape.top, self._measurement_unit)
            + unit_conversion(shape.height, self._measurement_unit),
        }

    def _generate_bounding_box(
        self, left: Length, top: Length, right: Length, bottom: Length
    ) -> dict:
        return {
            "left": unit_conversion(left, self._measurement_unit),
            "top": unit_conversion(top, self._measurement_unit),
            "right": unit_conversion(right, self._measurement_unit),
            "bottom": unit_conversion(bottom, self._measurement_unit),
        }

    def _try_split(self, shapes: list[Shape], direction: str) -> list[list[Shape]]:
        if not shapes:
            return []  # Return empty if no shapes to split

        # Define grid lines based on direction
        grid_lines = self._define_grid_lines(shapes, direction)
        shapes_number = len(shapes)
        for line in grid_lines:
            if direction == "horizontal":
                top, bottom = self._split_by_line(shapes, line, direction)

                if top and bottom and self._valid_split(top, bottom, shapes_number):
                    # print("\n---------------------------")
                    # print(f"Splitting at {line} {direction}")
                    # print(f"Top: {len(top)} shapes")
                    # print(f"Bottom: {len(bottom)} shapes")
                    # print("---------------------------")
                    return [top, bottom]
            elif direction == "vertical":
                left, right = self._split_by_line(shapes, line, direction)
                if left and right and self._valid_split(left, right, shapes_number):
                    # print("\n---------------------------")
                    # print(f"Splitting at {line} {direction}")
                    # print(f"Left: {len(left)} shapes")
                    # print(f"Right: {len(right)} shapes")
                    # print("---------------------------")
                    return [left, right]

        return []  # Return empty if no valid split found

    def _define_grid_lines(self, shapes: list[Shape], direction: str) -> list[float]:
        if not shapes:
            return []  # Return empty list if no shapes to define grid lines
        lines = []
        if direction == "horizontal":
            top_y_position = min(
                [unit_conversion(shape.top, self._measurement_unit) for shape in shapes]
            )
            bottom_y_position = max(
                [
                    unit_conversion(shape.top, self._measurement_unit)
                    + unit_conversion(shape.height, self._measurement_unit)
                    for shape in shapes
                ]
            )
            # print(f"Top: {top_y_position}, Bottom: {bottom_y_position}")
            intervals = [(top_y_position, bottom_y_position)]
            # print(f"Intervals: {intervals}")
            for shape in shapes:
                top = unit_conversion(shape.top, self._measurement_unit)
                height = unit_conversion(shape.height, self._measurement_unit)
                intervals = intervals_minus_interval(intervals, (top, top + height))
                # print(f"Intervals: {intervals}")
            # print(f"Intervals: {intervals}")
            for interval in intervals:
                line = (interval[0] + interval[1]) / 2
                # print(f"Line: {line}")
                lines.append(line)

            # print(f"Horizontal lines: {lines}")
            return sorted(lines)

        elif direction == "vertical":
            left_x_position = min(
                [
                    unit_conversion(shape.left, self._measurement_unit)
                    for shape in shapes
                ]
            )
            right_x_position = max(
                [
                    unit_conversion(shape.left, self._measurement_unit)
                    + unit_conversion(shape.width, self._measurement_unit)
                    for shape in shapes
                ]
            )
            intervals = [(left_x_position, right_x_position)]
            for shape in shapes:
                left = unit_conversion(shape.left, self._measurement_unit)
                width = unit_conversion(shape.width, self._measurement_unit)
                intervals = intervals_minus_interval(intervals, (left, left + width))

            for interval in intervals:
                line = (interval[0] + interval[1]) / 2
                lines.append(line)

            # print(f"Vertical lines: {lines}")
            return sorted(lines)

        else:
            raise ValueError(f"Invalid direction: {direction}")

    def _valid_split(
        self, group1: list[Shape], group2: list[Shape], shapes_number: int
    ) -> bool:
        if not group1 or not group2:
            return False

        for shape1 in group1:
            if shape1 in group2:
                return False

        # Ensure no shape is left out
        if len(group1) + len(group2) != shapes_number:
            return False

        return True

    def _split_by_line(self, shapes: list[Shape], line: float, direction: str) -> tuple:
        if direction == "horizontal":
            top = [
                shape
                for shape in shapes
                if unit_conversion(shape.top, self._measurement_unit)
                + unit_conversion(shape.height, self._measurement_unit)
                <= line
            ]
            bottom = [
                shape
                for shape in shapes
                if unit_conversion(shape.top, self._measurement_unit) >= line
            ]
            return top, bottom
        elif direction == "vertical":
            left = [
                shape
                for shape in shapes
                if unit_conversion(shape.left, self._measurement_unit)
                + unit_conversion(shape.width, self._measurement_unit)
                < line
            ]
            right = [
                shape
                for shape in shapes
                if unit_conversion(shape.left, self._measurement_unit) >= line
            ]
            return left, right
        else:
            raise ValueError(f"Invalid direction: {direction}")


class PowerPointSegmenter:
    def __init__(self, presentation: Presentation, measurement_unit: str = "pt"):
        self._presentation = presentation
        self._measurement_unit = measurement_unit
        self._slide_width = presentation.slide_width
        self._slide_height = presentation.slide_height

    def segment(self, slide_index: int) -> SegmentTreeNode | None:
        slide = self._presentation.slides[slide_index]
        shapes = slide.shapes
        if not shapes:
            return None
        return Segmenter(
            shapes, self._slide_width, self._slide_height, self._measurement_unit
        ).segment()

    def segment_all(self) -> dict:
        return {i: self.segment(i) for i in range(len(self._presentation.slides))}
